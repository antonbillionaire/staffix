"""
Staffix Sales Pitch Presentation Generator
Brand: #3B82F6 → #8B5CF6, Dark BG #0B0A1A, Font: Inter/Calibri
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_BG = RGBColor(0x0B, 0x0A, 0x1A)
DARK_CARD = RGBColor(0x12, 0x12, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
LIGHT_BLUE_BG = RGBColor(0x1E, 0x1B, 0x4B)

# Paths
LOGO_PATH = r"c:\Users\anton\Documents\GITHUB REPOSITORIUM\staffix\public\new-logo\A2_right_white_text.png"
ICON_PATH = r"c:\Users\anton\Documents\GITHUB REPOSITORIUM\staffix\public\new-logo\2_sfx_icon.png"
OUTPUT = r"c:\Users\anton\Documents\GITHUB REPOSITORIUM\staffix\Staffix_Pitch_RU.pptx"

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

FONT = "Calibri"  # Inter fallback - Calibri is universal


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p


def add_gradient_rect(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_card(slide, left, top, width, height, color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_logo_small(slide):
    """Add small logo to bottom-right corner"""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(6.8), height=Inches(0.45))


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Glow circles (decorative)
glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(0.5), Inches(4), Inches(4))
glow1.fill.solid()
glow1.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3E)
glow1.line.fill.background()

glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(3), Inches(4.5), Inches(4.5))
glow2.fill.solid()
glow2.fill.fore_color.rgb = RGBColor(0x1A, 0x15, 0x35)
glow2.line.fill.background()

# Logo centered
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(4.2), Inches(1.8), height=Inches(1.2))

# Tagline
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1),
             "AI-сотрудник для вашего бизнеса", 44, WHITE, True, PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.8),
             "Отвечает клиентам. Записывает на услуги. Работает 24/7.", 24, MUTED, False, PP_ALIGN.CENTER)

# URL
add_text_box(slide, Inches(5), Inches(6.2), Inches(3.3), Inches(0.5),
             "staffix.io", 18, BLUE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
             "Знакомо?", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

# Three problem cards
problems = [
    ("Клиент написал в 23:00", "Никто не ответил.\nУтром он уже у конкурента."),
    ("Администратор забыл\nперезвонить", "Потерян клиент\nна 500 000 ₸."),
    ("47 непрочитанных\nв WhatsApp", "А рабочий день\nуже закончился."),
]

for i, (title, desc) in enumerate(problems):
    x = Inches(0.8 + i * 4.1)
    card = add_card(slide, x, Inches(2), Inches(3.7), Inches(4))

    # Red X icon
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(1), Inches(0.6),
                 "✕", 36, RED, True)

    add_text_box(slide, x + Inches(0.3), Inches(3), Inches(3.1), Inches(1.2),
                 title, 22, WHITE, True)

    add_text_box(slide, x + Inches(0.3), Inches(4.2), Inches(3.1), Inches(1.5),
                 desc, 18, MUTED)

# Photo placeholder
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(6), Inches(0.5),
             "📷 Вставьте фото: уставший владелец бизнеса с телефоном", 12, RGBColor(0x64, 0x64, 0x80))


# ============================================================
# SLIDE 3 — Cost of Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
             "Сколько вы теряете?", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

# Big stats
stats = [
    ("68%", "клиентов уходят к конкуренту,\nесли не получили ответ в течение часа"),
    ("15-30", "клиентов в месяц теряет\nсредний бизнес из-за медленных ответов"),
    ("3 000 000 ₸", "упущенной выручки в год —\nминимальная оценка"),
]

for i, (num, desc) in enumerate(stats):
    y = Inches(2 + i * 1.7)
    card = add_card(slide, Inches(0.8), y, Inches(11.7), Inches(1.4))
    add_text_box(slide, Inches(1.2), y + Inches(0.15), Inches(3.5), Inches(1),
                 num, 48, BLUE, True)
    add_text_box(slide, Inches(5), y + Inches(0.25), Inches(7), Inches(1),
                 desc, 20, MUTED)


# ============================================================
# SLIDE 4 — Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Встречайте вашего AI-сотрудника", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

# Icon
if os.path.exists(ICON_PATH):
    slide.shapes.add_picture(ICON_PATH, Inches(0.8), Inches(2), height=Inches(1.5))

# Description
txBox = add_text_box(slide, Inches(2.8), Inches(2), Inches(9.5), Inches(3.5),
                     "Staffix — это AI, который знает ваш бизнес досконально:", 24, WHITE, True)
tf = txBox.text_frame
points = [
    "Ваши услуги, цены, акции",
    "График работы и свободные окна",
    "Как отвечать, в каком тоне",
    "FAQ и частые вопросы клиентов",
]
for pt in points:
    add_paragraph(tf, f"→  {pt}", 22, MUTED, space_before=8)

# Bottom tagline
card = add_card(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2))
add_text_box(slide, Inches(1.2), Inches(5.65), Inches(11), Inches(1),
             "Не увольняется  •  Не болеет  •  Не забывает  •  Работает 24/7",
             26, GREEN, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — How it Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "3 шага до AI-сотрудника", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

steps = [
    ("1", "Расскажите о бизнесе", "5 минут", "Услуги, цены, FAQ —\nAI запомнит всё"),
    ("2", "Подключите каналы", "2 минуты", "WhatsApp, Instagram,\nTelegram, Facebook"),
    ("3", "AI начинает работать", "Сразу", "Отвечает клиентам\nот имени вашего бизнеса"),
]

for i, (num, title, time, desc) in enumerate(steps):
    x = Inches(0.8 + i * 4.1)
    card = add_card(slide, x, Inches(2), Inches(3.7), Inches(4.5))

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(2.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.3), Inches(2.45), Inches(0.8), Inches(0.8),
                 num, 32, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.3), Inches(3.4), Inches(3.1), Inches(0.7),
                 title, 22, WHITE, True)

    add_text_box(slide, x + Inches(0.3), Inches(4.0), Inches(3.1), Inches(0.5),
                 time, 16, GREEN, True)

    add_text_box(slide, x + Inches(0.3), Inches(4.5), Inches(3.1), Inches(1.5),
                 desc, 18, MUTED)


# ============================================================
# SLIDE 6 — Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Давайте покажу прямо сейчас", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

# Demo scenario card
card = add_card(slide, Inches(1.5), Inches(2), Inches(10.3), Inches(4.5))

txBox = add_text_box(slide, Inches(2), Inches(2.3), Inches(9.3), Inches(0.5),
                     "Живая демонстрация:", 28, WHITE, True)
tf = txBox.text_frame

scenarios = [
    "Клиент спрашивает про услуги → AI отвечает с ценами",
    "Клиент хочет записаться → AI предлагает свободное время",
    "Клиент спрашивает адрес → AI даёт точный ответ",
    "Клиент пишет в нерабочее время → AI всё равно отвечает",
]
for sc in scenarios:
    add_paragraph(tf, f"✓  {sc}", 22, MUTED, space_before=12)

add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
             "📷 Переключитесь на телефон для живого демо", 16, RGBColor(0x64, 0x64, 0x80))


# ============================================================
# SLIDE 7 — Channels
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Один AI — все мессенджеры", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

channels = [
    ("WhatsApp", "#25D366", "Самый популярный\nв КЗ и УЗ"),
    ("Instagram", "#E1306C", "DM + комментарии\nк постам и рилс"),
    ("Telegram", "#0088CC", "Боты + группы\nдля бизнеса"),
    ("Facebook", "#1877F2", "Messenger\nдля страниц"),
]

channel_colors = [RGBColor(0x25, 0xD3, 0x66), RGBColor(0xE1, 0x30, 0x6C),
                  RGBColor(0x00, 0x88, 0xCC), RGBColor(0x18, 0x77, 0xF2)]

for i, (name, _, desc) in enumerate(channels):
    x = Inches(0.6 + i * 3.2)
    card = add_card(slide, x, Inches(2), Inches(2.9), Inches(3.5))

    # Channel color accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2), Inches(2.9), Pt(4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = channel_colors[i]
    accent.line.fill.background()

    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.6),
                 name, 28, WHITE, True)

    add_text_box(slide, x + Inches(0.2), Inches(3.3), Inches(2.5), Inches(1.5),
                 desc, 18, MUTED)

# Bottom note
card = add_card(slide, Inches(0.8), Inches(6), Inches(11.7), Inches(0.8))
add_text_box(slide, Inches(1.2), Inches(6.1), Inches(11), Inches(0.6),
             "Все переписки — в одном дашборде. Вы видите каждого клиента.",
             20, WHITE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8 — AI Capabilities
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Что умеет AI-сотрудник", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

capabilities = [
    ("💬", "Отвечает клиентам", "Мгновенно, 24/7, на языке клиента"),
    ("📅", "Записывает на услуги", "Проверяет свободные окна в расписании"),
    ("🎯", "Квалифицирует лидов", "cold → warm → hot → клиент"),
    ("📋", "Знает ваш каталог", "Цены, акции, описания товаров"),
    ("🌐", "4 языка", "Русский, английский, казахский, узбекский"),
    ("📊", "Ведёт аналитику", "Обращения, конверсия, каналы"),
]

for i, (icon, title, desc) in enumerate(capabilities):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2 + row * 2.5)

    card = add_card(slide, x, y, Inches(3.7), Inches(2.2))

    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6),
                 icon, 28, WHITE)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.3), Inches(0.5),
                 title, 20, WHITE, True)

    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), Inches(3.3), Inches(0.6),
                 desc, 16, MUTED)


# ============================================================
# SLIDE 9 — Target Audience
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Для какого бизнеса подходит", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

businesses = [
    ("💈", "Салоны красоты", "Запись, прайс, напоминания"),
    ("🍽️", "Рестораны и кафе", "Меню, бронь столов, доставка"),
    ("🏥", "Клиники", "Запись к врачу, FAQ по услугам"),
    ("🏋️", "Фитнес-клубы", "Расписание, абонементы"),
    ("🛒", "Интернет-магазины", "Каталог, статус заказа"),
    ("🔧", "Сервисные компании", "Заявки, статус ремонта"),
]

for i, (icon, title, desc) in enumerate(businesses):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2 + row * 2.5)

    card = add_card(slide, x, y, Inches(3.7), Inches(2.2))

    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.8), Inches(0.8),
                 icon, 36, WHITE)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.3), Inches(0.5),
                 title, 22, WHITE, True)

    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), Inches(3.3), Inches(0.6),
                 desc, 16, MUTED)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
             "📷 Вставьте фото: довольные владельцы бизнесов / интерьеры салонов и ресторанов", 12, RGBColor(0x64, 0x64, 0x80))


# ============================================================
# SLIDE 10 — Comparison Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Администратор vs AI-сотрудник", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

# Table headers
header_card = add_card(slide, Inches(0.8), Inches(1.8), Inches(3.5), Inches(0.7), DARK_CARD)
add_text_box(slide, Inches(0.8), Inches(1.85), Inches(3.5), Inches(0.6),
             "", 18, MUTED, False, PP_ALIGN.CENTER)

admin_header = add_card(slide, Inches(4.5), Inches(1.8), Inches(4), Inches(0.7), DARK_CARD)
add_text_box(slide, Inches(4.5), Inches(1.85), Inches(4), Inches(0.6),
             "👤 Администратор", 20, MUTED, True, PP_ALIGN.CENTER)

ai_header = add_card(slide, Inches(8.7), Inches(1.8), Inches(4), Inches(0.7))
ai_header.fill.solid()
ai_header.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3E)
add_text_box(slide, Inches(8.7), Inches(1.85), Inches(4), Inches(0.6),
             "🤖 AI Staffix", 20, BLUE, True, PP_ALIGN.CENTER)

rows = [
    ("Стоимость", "150-250K ₸/мес", "от 5 000 ₸/мес"),
    ("График", "8 часов", "24/7/365"),
    ("Скорость ответа", "5-30 минут", "3 секунды"),
    ("Забывает?", "Да", "Никогда"),
    ("Болеет?", "Да", "Нет"),
    ("Масштаб", "1 чат за раз", "100+ одновременно"),
    ("Языки", "1-2", "4"),
]

for i, (label, admin_val, ai_val) in enumerate(rows):
    y = Inches(2.7 + i * 0.65)
    stripe_color = DARK_CARD if i % 2 == 0 else DARK_BG

    add_card(slide, Inches(0.8), y, Inches(3.5), Inches(0.6), stripe_color)
    add_text_box(slide, Inches(1), y + Inches(0.05), Inches(3.3), Inches(0.5),
                 label, 17, WHITE, False)

    add_card(slide, Inches(4.5), y, Inches(4), Inches(0.6), stripe_color)
    add_text_box(slide, Inches(4.5), y + Inches(0.05), Inches(0.5), Inches(0.5),
                 "✕", 17, RED, True)
    add_text_box(slide, Inches(5), y + Inches(0.05), Inches(3.3), Inches(0.5),
                 admin_val, 17, MUTED)

    add_card(slide, Inches(8.7), y, Inches(4), Inches(0.6), stripe_color)
    add_text_box(slide, Inches(8.7), y + Inches(0.05), Inches(0.5), Inches(0.5),
                 "✓", 17, GREEN, True)
    add_text_box(slide, Inches(9.2), y + Inches(0.05), Inches(3.3), Inches(0.5),
                 ai_val, 17, GREEN)


# ============================================================
# SLIDE 11 — Pricing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Простые тарифы", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

plans = [
    ("Start", "5 000 ₸", "/мес", ["100 сообщений", "1 канал", "Базовая аналитика", "Email поддержка"]),
    ("Pro", "15 000 ₸", "/мес", ["1 000 сообщений", "Все каналы", "Полная аналитика", "Квалификация лидов", "Приоритетная поддержка"]),
    ("Business", "35 000 ₸", "/мес", ["Безлимит сообщений", "Все каналы", "API доступ", "Персональная настройка", "Выделенный менеджер"]),
]

for i, (name, price, period, features) in enumerate(plans):
    x = Inches(0.8 + i * 4.1)
    is_popular = (i == 1)

    if is_popular:
        # Gradient-ish border effect
        border = add_card(slide, x - Inches(0.05), Inches(1.95), Inches(3.8), Inches(5.1), BLUE)
        card = add_card(slide, x, Inches(2), Inches(3.7), Inches(5))
    else:
        card = add_card(slide, x, Inches(2), Inches(3.7), Inches(5))

    # Badge for popular
    if is_popular:
        badge = add_card(slide, x + Inches(0.8), Inches(1.8), Inches(2.1), Inches(0.45), BLUE)
        add_text_box(slide, x + Inches(0.8), Inches(1.82), Inches(2.1), Inches(0.4),
                     "РЕКОМЕНДУЕМ", 13, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.3), Inches(2.4), Inches(3.1), Inches(0.5),
                 name, 26, BLUE if is_popular else MUTED, True)

    add_text_box(slide, x + Inches(0.3), Inches(3), Inches(3.1), Inches(0.7),
                 price, 40, WHITE, True)

    add_text_box(slide, x + Inches(0.3) + Inches(2.2), Inches(3.2), Inches(1), Inches(0.5),
                 period, 16, MUTED)

    for j, feat in enumerate(features):
        add_text_box(slide, x + Inches(0.3), Inches(3.9 + j * 0.5), Inches(3.1), Inches(0.4),
                     f"✓  {feat}", 16, WHITE if is_popular else MUTED)


# ============================================================
# SLIDE 12 — Social Proof
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Уже работает", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

# Quote card
card = add_card(slide, Inches(0.8), Inches(2), Inches(11.7), Inches(2.2))
add_text_box(slide, Inches(1.3), Inches(2.2), Inches(0.5), Inches(0.5),
             "\"", 72, BLUE, True)
add_text_box(slide, Inches(1.8), Inches(2.5), Inches(10), Inches(1),
             "За первую неделю AI обработал 200+ обращений.\nМы перестали терять ночных клиентов.",
             24, WHITE, False)
add_text_box(slide, Inches(1.8), Inches(3.5), Inches(10), Inches(0.5),
             "— Салон красоты, Алматы", 18, MUTED)

# Stats
stats_data = [
    ("50+", "бизнесов\nподключено"),
    ("10 000+", "обработанных\nсообщений"),
    ("4 страны", "KZ, UZ,\nRU, KR"),
]

for i, (num, label) in enumerate(stats_data):
    x = Inches(0.8 + i * 4.1)
    card = add_card(slide, x, Inches(4.8), Inches(3.7), Inches(2))
    add_text_box(slide, x + Inches(0.3), Inches(5), Inches(3.1), Inches(0.8),
                 num, 40, BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(5.8), Inches(3.1), Inches(0.8),
                 label, 18, MUTED, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(7), Inches(8), Inches(0.3),
             "* Обновите цифры на актуальные", 12, RGBColor(0x64, 0x64, 0x80))


# ============================================================
# SLIDE 13 — CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(0.8),
             "Начните получать клиентов уже сегодня", 44, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

# Steps
cta_steps = [
    ("1", "Регистрация", "1 минута"),
    ("2", "Настройка бизнеса", "5 минут"),
    ("3", "AI отвечает клиентам", "Сразу"),
]

for i, (num, title, time) in enumerate(cta_steps):
    x = Inches(0.8 + i * 4.1)
    card = add_card(slide, x, Inches(2), Inches(3.7), Inches(2))

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(2.3), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(0.6), Inches(0.6),
                 num, 24, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.1), Inches(2.3), Inches(2.3), Inches(0.4),
                 title, 22, WHITE, True)
    add_text_box(slide, x + Inches(1.1), Inches(2.8), Inches(2.3), Inches(0.4),
                 time, 18, GREEN)

# Free trial card
card = add_card(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(2.2))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3E)

add_text_box(slide, Inches(2.5), Inches(4.7), Inches(8.3), Inches(0.7),
             "Первые 50 сообщений — бесплатно", 32, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.5), Inches(5.4), Inches(8.3), Inches(0.5),
             "Никаких обязательств. Отмена в 1 клик.", 20, MUTED, False, PP_ALIGN.CENTER)

# CTA button
btn = add_card(slide, Inches(4.5), Inches(6), Inches(4.3), Inches(0.8), BLUE)
add_text_box(slide, Inches(4.5), Inches(6.05), Inches(4.3), Inches(0.7),
             "Попробовать бесплатно →", 22, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(4.5), Inches(6.9), Inches(4.3), Inches(0.4),
             "staffix.io", 18, BLUE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 14 — Q&A / Final
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Decorative glows (subtle)
glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(1), Inches(3.5), Inches(3.5))
glow1.fill.solid()
glow1.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3E)
glow1.line.fill.background()

glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(3.5), Inches(4), Inches(4))
glow2.fill.solid()
glow2.fill.fore_color.rgb = RGBColor(0x1A, 0x15, 0x35)
glow2.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
             "Вопросы?", 60, WHITE, True, PP_ALIGN.CENTER)

# Logo
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(4.5), Inches(3.2), height=Inches(0.9))

# Contacts
contacts = [
    "🌐  staffix.io",
    "📸  @staffixio",
    "💬  Telegram: @staffix_support",
]
for i, contact in enumerate(contacts):
    add_text_box(slide, Inches(3.5), Inches(4.8 + i * 0.55), Inches(6.3), Inches(0.5),
                 contact, 22, MUTED, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.5),
             "Спасибо за внимание!", 24, WHITE, True, PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
