import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

prs = Presentation('Staffix_Pitch_RU_08.04.26.pptx')

# Full paragraph text -> English translation
T = {
    # Slide 1
    "AI-сотрудник для вашего бизнеса, которого вы так искали!": "The AI Employee your business has been looking for!",
    # Slide 2
    "Знакомо?": "Sound familiar?",
    "89% компаний МСБ нанимают недостаточно квалифицированных сотрудников": "89% of SMBs hire underqualified employees",
    "Среднее время обучения нового сотрудника- 40 часов. Продуктивность- через 3-4 месяца": "Average new hire training: 40 hours. Full productivity: 3\u20134 months",
    "60% персонала находят другую работу в течение первых 3 месяцев": "60% of staff find another job within the first 3 months",
    # Slide 3
    "А знали ли вы?": "Did you know?",
    "Клиент написал \u2013 вы ему ответили, а он уже купил/заказал у другой компании": "Customer reached out \u2014 by the time you replied, they already bought from a competitor",
    "Клиент написал – вы ему ответили, а он уже купил/заказал у другой компании": "Customer reached out \u2014 by the time you replied, they already bought from a competitor",
    "78% лидов покупают у того, кто ответил первым": "78% of leads buy from whoever responds first",
    "Клиент с бОльшей охотой напишет вам, чем позвонит": "Customers prefer messaging over calling",
    "68% клиентов предпочитают мессенджеры- именно из-за скорости": "68% of customers prefer messengers \u2014 because of speed",
    "Ответ в течение первых 5 минут- в 100 раз выше шанс конверсии": "Responding within 5 minutes \u2014 100x higher conversion chance",
    "85 % клиентов больше не напишут вам, если не получат ответ": "85% of customers won\u2019t write again if they don\u2019t get a reply",
    "82% из них хотят ответ в течение 10 минут": "82% of them want a reply within 10 minutes",
    # Slide 4
    "А вы даже и незнали": "And you didn\u2019t even know",
    "Средний бизнес отвечает на лида в течение 47 часов": "Average business responds to a lead in 47 hours",
    "Только 20% МСБ отвечают в течение первого часа": "Only 20% of SMBs respond within the first hour",
    "60% звонков и сообщений остаются без ответа": "60% of calls and messages go unanswered",
    "12.15 дол- средняя стоимость 1 пропущенного звонка": "$12.15 \u2014 average cost of 1 missed call",
    "1 пропущенный звонок в день это 4.400 дол потерь в год": "1 missed call per day = $4,400 in losses per year",
    # Slide 5
    "Встречайте вашего AI-сотрудника": "Meet your AI Employee",
    "Staffix:": "Staffix:",
    "Учится за 5 минут и знает Ваш бизнес досконально": "Learns in 5 minutes and knows your business inside out",
    "Отвечает Вашим клиентам в течение 3-5 !! секунд": "Responds to your customers in 3\u20135 !! seconds",
    "Может полностью автоматизировать работу фронт-офиса": "Can fully automate your front-office operations",
    "Любит и умеет продавать, работать с возражениями, квалифицировать лидов, и вести учет": "Sells, handles objections, qualifies leads, and keeps records",
    "И многое другое": "And much more",
    "Не увольняется  \u2022  Не болеет  \u2022  Не забывает  \u2022  Работает 24/7": "Never quits  \u2022  Never sick  \u2022  Never forgets  \u2022  Works 24/7",
    # Slide 6
    "Что умеет AI-сотрудник": "What the AI Employee can do",
    "Общается с клиентами: консультации, продажи, работа с возражениями": "Communicates with customers: consultations, sales, objection handling",
    "Мгновенно, 24/7, на языке клиента (Узбекский, Казахский, Русский, Английский)": "Instantly, 24/7, in customer\u2019s language (Uzbek, Kazakh, Russian, English)",
    "Записывает на услуги/ принимает заказы": "Books appointments / takes orders",
    "Проверяет расписание, наличие товара, оформляет заказ, отправляет уведомления, ведет учет": "Checks schedules, stock, places orders, sends notifications, keeps records",
    "Квалифицирует лидов": "Qualifies leads",
    "И ведет их до момента покупки cold \u2192 warm \u2192 hot \u2192 клиент": "And nurtures them to purchase: cold \u2192 warm \u2192 hot \u2192 customer",
    "+ Автоматизация: напоминания, реанимация, сбор отзывов, рассылки, акции": "+ Automation: reminders, reactivation, review collection, broadcasts, promotions",
    "Знает ваш бизнес": "Knows your business",
    "Мастеров, клиентов, каталог, цены, описание товаров и услуг, политики, программы лояльности": "Staff, customers, catalog, prices, descriptions, policies, loyalty programs",
    "Мультиканальный": "Multichannel",
    "Одновременно общается в 4 каналах: Telegram, WhatsApp, Instagram, Facebook.": "Communicates across 4 channels: Telegram, WhatsApp, Instagram, Facebook.",
    "Ведёт аналитику": "Tracks analytics",
    "Обращения, конверсия, каналы, продажи, заработок": "Inquiries, conversion, channels, sales, revenue",
    # Slide 7
    "3 шага до AI-сотрудника": "3 steps to your AI Employee",
    "Зарегистрируйтесь на Staffix.io": "Sign up at Staffix.io",
    "2 минуты": "2 minutes",
    "Регистрация через имейл, либо через Google авторизацию": "Register via email or Google sign-in",
    "Настройте и подключите AI сотрудника": "Set up and connect your AI Employee",
    "5 минут": "5 minutes",
    "Загрузите Базу знаний": "Upload your Knowledge Base",
    "Подключите Telegram, WA, IG, FB": "Connect Telegram, WA, IG, FB",
    "Настройте команду, товары/услуги": "Set up your team, products/services",
    "AI сотрудник готов к работе": "AI Employee is ready to work",
    "Сразу": "Instantly",
    "Отвечает клиентамот имени вашего бизнеса во всех настроенных каналах": "Responds to customers on behalf of your business across all connected channels",
    # Slide 8
    "Попробуйте прямо сейчас": "Try it right now",
    # Slide 9
    "Один AI \u2014 все мессенджеры": "One AI \u2014 all messengers",
    "Самый популярныйв КЗ": "Most popular in KZ",
    "DM + комментариик постам и рилс": "DM + comments on posts and reels",
    "Самый популярный в Узбекистане": "Most popular in Uzbekistan",
    "Messengerдля страниц": "Messenger for pages",
    "Все переписки \u2014 в одном дашборде. Вы видите каждого клиента.": "All conversations in one dashboard. You see every customer.",
    # Slide 10
    "Для какого бизнеса подходит": "What businesses it works for",
    "СЕРВИСНЫЕ БИЗНЕСЫ": "SERVICE BUSINESSES",
    "ТОРГОВЫЕ БИЗНЕСЫ": "RETAIL BUSINESSES",
    "Салоны красоты": "Beauty salons",
    "Барбершопы": "Barbershops",
    "Клиники": "Clinics",
    "Диагностические центры": "Diagnostic centers",
    "СПА": "SPA",
    "Фитнес": "Fitness",
    "Клининг": "Cleaning",
    "Уход за животными": "Pet care",
    "Образовательные центры": "Education centers",
    "Профессиональные услуги": "Professional services",
    "Магазины цветов": "Flower shops",
    "Онлайн магазины": "Online stores",
    "Оптовые магазины": "Wholesale stores",
    "Розничные магазины": "Retail stores",
    "Торговые точки с доставкой": "Stores with delivery",
    # Slide 11
    "Администратор vs AI-сотрудник": "Receptionist vs AI Employee",
    "Администратор": "Receptionist",
    "AI Staffix": "AI Staffix",
    "Стоимость": "Cost",
    "300-500 дол/мес": "$300\u2013500/mo",
    "от 20 дол/мес": "from $20/mo",
    "График": "Schedule",
    "8-10 часов": "8\u201310 hours",
    "Скорость ответа": "Response time",
    "5-30 минут": "5\u201330 minutes",
    "3 секунды": "3 seconds",
    "Забывает?": "Forgets?",
    "Да": "Yes",
    "Никогда": "Never",
    "Болеет?": "Gets sick?",
    "Нет": "No",
    "Масштаб": "Scale",
    "1 чат за раз": "1 chat at a time",
    "100+ одновременно": "100+ simultaneously",
    "Языки": "Languages",
    # Slide 12
    "Простые тарифы": "Simple pricing",
    # Slide 13
    "Сколько зарабатывает салон с Staffix?": "How much does a salon earn with Staffix?",
    "Салон красоты  \u2022  3 мастера  \u2022  средний чек 200 000 сум": "Beauty salon  \u2022  3 stylists  \u2022  avg check 200,000 UZS",
    "Дополнительная выручка с AI-сотрудником:": "Additional revenue with AI Employee:",
    "Быстрые ответы \u2192 больше записей": "Fast replies \u2192 more bookings",
    "Конверсия обращений: 60% \u2192 85%  \u2022  +1 записи/день": "Inquiry conversion: 60% \u2192 85%  \u2022  +1 booking/day",
    "+6 000 000 сум/мес": "+6,000,000 UZS/mo",
    "Напоминания \u2192 меньше пустых кресел": "Reminders \u2192 fewer empty chairs",
    "Неявки снижаются на 50%  \u2022  15 спасённых записей/мес": "No-shows reduced by 50%  \u2022  15 saved bookings/mo",
    "+3 000 000 сум/мес": "+3,000,000 UZS/mo",
    "Реактивация \u2192 возврат клиентов": "Reactivation \u2192 returning customers",
    "20 неактивных клиентов  \u2022  20% возвращаются": "20 inactive clients  \u2022  20% return",
    "+800 000 сум/мес": "+800,000 UZS/mo",
    "Отзывы \u2192 рост рейтинга": "Reviews \u2192 rating growth",
    "Рейтинг 4.2 \u2192 4.7 в Google/2GIS  \u2022  +10% новых обращений": "Rating 4.2 \u2192 4.7 on Google/2GIS  \u2022  +10% new inquiries",
    "+новые клиенты": "+new customers",
    "ИТОГО": "TOTAL",
    "+9 800 000": "+9,800,000",
    "сум/мес доп. выручки": "UZS/mo additional revenue",
    "Стоимость Staffix:": "Staffix cost:",
    "550 000 сум/мес": "550,000 UZS/mo",
    "Окупаемость:": "Payback:",
    "за 1 день": "in 1 day",
    "Staffix стоит чуть больше одного среднего чека. А приносит \u2014  миллионы недополученной прибыли.": "Staffix costs slightly more than one average check. But brings back millions in lost revenue.",
    "staffix.io  \u2022  14 дней бесплатно": "staffix.io  \u2022  14 days free",
    # Slide 14
    "Сколько зарабатывает клиника с Staffix?": "How much does a clinic earn with Staffix?",
    "Клиника / Стоматология  \u2022  3 врача  \u2022  средний чек 500 000 сум": "Clinic / Dental  \u2022  3 doctors  \u2022  avg check 500,000 UZS",
    "Быстрые ответы \u2192 больше пациентов": "Fast replies \u2192 more patients",
    "Конверсия обращений: 50% \u2192 75%  \u2022  +2-3 записей/день": "Inquiry conversion: 50% \u2192 75%  \u2022  +2\u20133 bookings/day",
    "+30 000 000 сум/мес": "+30,000,000 UZS/mo",
    "Напоминания \u2192 меньше пустых слотов": "Reminders \u2192 fewer empty slots",
    "Неявки снижаются на 55%  \u2022  11 спасённых приёма/мес": "No-shows reduced by 55%  \u2022  11 saved appointments/mo",
    "+5 500 000 сум/мес": "+5,500,000 UZS/mo",
    "Реактивация \u2192 возврат пациентов": "Reactivation \u2192 returning patients",
    "40 неактивных пациентов  \u2022  'Пора на осмотр'  \u2022  15% возвращаются": "40 inactive patients  \u2022  'Time for a checkup'  \u2022  15% return",
    "Отзывы \u2192 доверие и поток": "Reviews \u2192 trust and traffic",
    "Рейтинг 4.3 \u2192 4.8  \u2022  Для клиники рейтинг = решение о визите": "Rating 4.3 \u2192 4.8  \u2022  For a clinic, rating = visit decision",
    "+новые пациенты": "+new patients",
    "+38 500 000": "+38,500,000",
    "за 1 приём": "in 1 appointment",
    "Один пропущенный пациент = 500 000 сум. Staffix стоит как один приём.": "One missed patient = 500,000 UZS. Staffix costs as much as one appointment.",
    # Slide 15
    "Сколько зарабатывает магазин с Staffix?": "How much does a store earn with Staffix?",
    "Магазин / Доставка  \u2022  средний чек 300 000 сум": "Store / Delivery  \u2022  avg check 300,000 UZS",
    "Быстрые ответы \u2192 больше заказов": "Fast replies \u2192 more orders",
    "Конверсия: 40% \u2192 65%  \u2022  +3 заказов/день": "Conversion: 40% \u2192 65%  \u2022  +3 orders/day",
    "+27 000 000 сум/мес": "+27,000,000 UZS/mo",
    "Допродажи \u2192 рост среднего чека": "Upsells \u2192 higher average check",
    "'К букету часто берут открытку и конфеты'  \u2022  +20% к чеку": "'Customers often add a card and chocolates'  \u2022  +20% to check",
    "+4 000 000 сум/мес": "+4,000,000 UZS/mo",
    "Реактивация \u2192 повторные заказы": "Reactivation \u2192 repeat orders",
    "50 неактивных клиентов  \u2022  'Новая коллекция! Скидка 10%'": "50 inactive clients  \u2022  'New collection! 10% off'",
    "Рассылки \u2192 акции по сегментам": "Broadcasts \u2192 promotions by segment",
    "VIP \u2014 эксклюзив  \u2022  Новые \u2014 скидка  \u2022  Неактивные \u2014 возврат": "VIP \u2014 exclusive  \u2022  New \u2014 discount  \u2022  Inactive \u2014 win-back",
    "+37 000 000": "+37,000,000",
    "1 160 000 сум/мес": "1,160,000 UZS/mo",
    "за 4 заказа": "in 4 orders",
    "AI-сотрудник консультирует, продаёт и допродаёт. Вы собираете и отправляете.": "AI Employee consults, sells and upsells. You pack and ship.",
    # Slide 16
    "Начните получать клиентов уже сегодня": "Start getting customers today",
    "Первые 100 сообщений \u2014 бесплатно": "First 100 messages \u2014 free",
    "14 дней. Никаких обязательств.": "14 days. No commitments.",
    "Попробуйте бесплатно": "Try for free",
    "прямо сейчас": "right now",
    # Slide 17
    "Спасибо за внимание!": "Thank you!",
    # Missed paragraphs (with arrows and merged words)
    "\u2192  Учится за 5 минут и знает Ваш бизнес досконально": "\u2192  Learns in 5 minutes and knows your business inside out",
    "\u2192  Отвечает Вашим клиентам в течение 3-5 !! секунд": "\u2192  Responds to your customers in 3\u20135 !! seconds",
    "\u2192  Может полностью автоматизировать работу фронт-офиса": "\u2192  Can fully automate your front-office operations",
    "\u2192  Любит и умеет продавать, работать с возражениями, квалифицировать лидов, и вести учет": "\u2192  Sells, handles objections, qualifies leads, and keeps records",
    "\u2192 И многое другое": "\u2192 And much more",
    "\u2192  \u0423\u0447\u0438\u0442\u0441\u044f \u0437\u0430 5 \u043c\u0438\u043d\u0443\u0442 \u0438 \u0437\u043d\u0430\u0435\u0442 \u0412\u0430\u0448 \u0431\u0438\u0437\u043d\u0435\u0441 \u0434\u043e\u0441\u043a\u043e\u043d\u0430\u043b\u044c\u043d\u043e": "\u2192  Learns in 5 minutes and knows your business inside out",
    "Отвечает клиентамот имени вашего бизнеса во всех настроенных каналах": "Responds to customers on behalf of your business across all connected channels",
    "Отвечает клиентам от имени вашего бизнеса во всех настроенных каналах": "Responds to customers on behalf of your business across all connected channels",
    "Самый популярныйв КЗ": "Most popular in KZ",
    "Самый популярный в КЗ": "Most popular in KZ",
    "DM + комментариик постам и рилс": "DM + comments on posts and reels",
    "DM + комментарии к постам и рилс": "DM + comments on posts and reels",
    "Messengerдля страниц": "Messenger for pages",
    "Messenger для страниц": "Messenger for pages",
    "\U0001f464 Администратор": "\U0001f464 Receptionist",
    "→  Учится за 5 минут и знает Ваш бизнес досконально": "→  Learns in 5 minutes and knows your business inside out",
    "→  Отвечает Вашим клиентам в течение 3-5 !! секунд": "→  Responds to your customers in 3\u20135 !! seconds",
    "→  Может полностью автоматизировать работу фронт-офиса": "→  Can fully automate your front-office operations",
    "→  Любит и умеет продавать, работать с возражениями, квалифицировать лидов, и вести учет": "→  Sells, handles objections, qualifies leads, and keeps records",
    "→ И многое другое": "→ And much more",
}

count = 0
missed = []

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                # Get full paragraph text
                full_text = para.text.strip()
                if not full_text:
                    continue

                # Normalize vertical tabs for matching
                full_text_norm = full_text.replace('\x0b', ' ')

                # Try exact match, then normalized (without vertical tabs)
                en_text = T.get(full_text) or T.get(full_text_norm)

                if en_text:
                    if para.runs:
                        para.runs[0].text = en_text
                        for run in para.runs[1:]:
                            run.text = ""
                        count += 1
                else:
                    # Check if it's a known untranslatable (numbers, symbols, etc.)
                    if full_text not in ["staffix.io", "Staffix:", "1", "2", "3", "4", "1-2",
                                         "24/7/365", "AI Staffix", "WhatsApp", "Instagram",
                                         "Telegram", "Facebook", "\u2715",
                                         "\u2713", "\u2612", "\u2611",
                                         "\u2716", "Staffix AI", "staffix.io",
                                         "@staffix.io", "Telegram: Staffix Client Manager",
                                         "Whatsapp: +8210-9869-6996",
                                         "\u250c\u2500\u2500\u2500\u2500"] and len(full_text) > 2:
                        # Check if it contains Cyrillic
                        has_cyrillic = any('\u0400' <= c <= '\u04ff' for c in full_text)
                        if has_cyrillic:
                            missed.append(full_text)

prs.save('Staffix_Pitch_EN_08.04.26.pptx')
print(f'Translated {count} paragraphs')
if missed:
    print(f'\nMissed {len(missed)} paragraphs with Cyrillic:')
    for m in missed:
        print(f'  "{m}"')
