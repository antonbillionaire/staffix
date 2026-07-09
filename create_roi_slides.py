"""Create 3 ROI slides for Staffix presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x0B, 0x0A, 0x1A)
DARK_CARD = RGBColor(0x12, 0x12, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
LIGHT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
GOLD = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def add_card(slide, left, top, width, height, color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_roi_line(slide, left, top, icon, title, desc, amount, amount_color=LIGHT_GREEN):
    """Add a ROI line item with icon, description and amount."""
    add_text(slide, left, top, 0.5, 0.4, icon, size=20, align=PP_ALIGN.CENTER)
    add_text(slide, left + 0.5, top - 0.05, 4.5, 0.35, title, size=16, bold=True)
    add_text(slide, left + 0.5, top + 0.3, 4.5, 0.3, desc, size=12, color=MUTED)
    add_text(slide, left + 5.2, top + 0.05, 2.5, 0.35, amount, size=18, bold=True, color=amount_color, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: SALON
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Title
add_text(slide, 0.8, 0.3, 12, 0.6, "Сколько зарабатывает салон с Staffix?", size=32, bold=True)
add_text(slide, 0.8, 0.85, 6, 0.4, "Салон красоты  •  3 мастера  •  средний чек 200 000 сум", size=14, color=MUTED)

# Left card - ROI items
add_card(slide, 0.5, 1.5, 8, 4.8)

add_text(slide, 1.0, 1.7, 7, 0.4, "Дополнительная выручка с AI-сотрудником:", size=16, color=BLUE, bold=True)

add_roi_line(slide, 1.0, 2.3,
    "⚡", "Быстрые ответы → больше записей",
    "Конверсия обращений: 60% → 85%  •  +4 записи/день",
    "+20 000 000 сум/мес")

add_roi_line(slide, 1.0, 3.1,
    "🔔", "Напоминания → меньше пустых кресел",
    "Неявки снижаются на 50%  •  15 спасённых записей/мес",
    "+3 000 000 сум/мес")

add_roi_line(slide, 1.0, 3.9,
    "🔄", "Реактивация → возврат клиентов",
    "50 неактивных клиентов  •  20% возвращаются",
    "+2 000 000 сум/мес")

add_roi_line(slide, 1.0, 4.7,
    "⭐", "Отзывы → рост рейтинга",
    "Рейтинг 4.2 → 4.7 в Google/2GIS  •  +10% новых обращений",
    "+новые клиенты")

# Divider line
add_text(slide, 1.0, 5.5, 7.2, 0.05, "─" * 80, size=8, color=MUTED)

# Right card - Summary
add_card(slide, 9, 1.5, 3.8, 4.8)

add_text(slide, 9.4, 1.8, 3.2, 0.4, "ИТОГО", size=14, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 9.4, 2.3, 3.2, 0.5, "+25 000 000", size=36, bold=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 2.9, 3.2, 0.3, "сум/мес доп. выручки", size=13, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 9.4, 3.6, 3.2, 0.3, "Стоимость Staffix:", size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 3.9, 3.2, 0.4, "260 000 сум/мес", size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text(slide, 9.4, 4.6, 3.2, 0.3, "Окупаемость:", size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 4.9, 3.2, 0.5, "за 1 день", size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Bottom
add_text(slide, 0.8, 6.6, 12, 0.4, "Staffix стоит меньше одного среднего чека. А приносит — десятки миллионов.", size=16, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 0.8, 7.0, 12, 0.3, "staffix.io  •  14 дней бесплатно", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: CLINIC
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "Сколько зарабатывает клиника с Staffix?", size=32, bold=True)
add_text(slide, 0.8, 0.85, 6, 0.4, "Клиника / Стоматология  •  3 врача  •  средний чек 500 000 сум", size=14, color=MUTED)

add_card(slide, 0.5, 1.5, 8, 4.8)

add_text(slide, 1.0, 1.7, 7, 0.4, "Дополнительная выручка с AI-сотрудником:", size=16, color=BLUE, bold=True)

add_roi_line(slide, 1.0, 2.3,
    "⚡", "Быстрые ответы → больше пациентов",
    "Конверсия обращений: 50% → 75%  •  +5 записей/день",
    "+62 500 000 сум/мес")

add_roi_line(slide, 1.0, 3.1,
    "🔔", "Напоминания → меньше пустых слотов",
    "Неявки снижаются на 55%  •  22 спасённых приёма/мес",
    "+11 000 000 сум/мес")

add_roi_line(slide, 1.0, 3.9,
    "🔄", "Реактивация → возврат пациентов",
    "80 неактивных пациентов  •  'Пора на осмотр'  •  15% возвращаются",
    "+6 000 000 сум/мес")

add_roi_line(slide, 1.0, 4.7,
    "⭐", "Отзывы → доверие и поток",
    "Рейтинг 4.3 → 4.8  •  Для клиники рейтинг = решение о визите",
    "+новые пациенты")

add_text(slide, 1.0, 5.5, 7.2, 0.05, "─" * 80, size=8, color=MUTED)

add_card(slide, 9, 1.5, 3.8, 4.8)

add_text(slide, 9.4, 1.8, 3.2, 0.4, "ИТОГО", size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 2.3, 3.2, 0.5, "+79 500 000", size=34, bold=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 2.9, 3.2, 0.3, "сум/мес доп. выручки", size=13, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 9.4, 3.6, 3.2, 0.3, "Стоимость Staffix:", size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 3.9, 3.2, 0.4, "520 000 сум/мес", size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text(slide, 9.4, 4.6, 3.2, 0.3, "Окупаемость:", size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 4.9, 3.2, 0.5, "за 1 приём", size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text(slide, 0.8, 6.6, 12, 0.4, "Один пропущенный пациент = 500 000 сум. Staffix стоит как один приём.", size=16, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 0.8, 7.0, 12, 0.3, "staffix.io  •  14 дней бесплатно", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: STORE / DELIVERY
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "Сколько зарабатывает магазин с Staffix?", size=32, bold=True)
add_text(slide, 0.8, 0.85, 6, 0.4, "Магазин / Доставка  •  средний чек 300 000 сум", size=14, color=MUTED)

add_card(slide, 0.5, 1.5, 8, 4.8)

add_text(slide, 1.0, 1.7, 7, 0.4, "Дополнительная выручка с AI-сотрудником:", size=16, color=BLUE, bold=True)

add_roi_line(slide, 1.0, 2.3,
    "⚡", "Быстрые ответы → больше заказов",
    "Конверсия: 40% → 65%  •  +6 заказов/день",
    "+45 000 000 сум/мес")

add_roi_line(slide, 1.0, 3.1,
    "🛒", "Допродажи → рост среднего чека",
    "'К букету часто берут открытку и конфеты'  •  +20% к чеку",
    "+9 000 000 сум/мес")

add_roi_line(slide, 1.0, 3.9,
    "🔄", "Реактивация → повторные заказы",
    "100 неактивных клиентов  •  'Новая коллекция! Скидка 10%'",
    "+6 000 000 сум/мес")

add_roi_line(slide, 1.0, 4.7,
    "📢", "Рассылки → акции по сегментам",
    "VIP — эксклюзив  •  Новые — скидка  •  Неактивные — возврат",
    "+3 000 000 сум/мес")

add_text(slide, 1.0, 5.5, 7.2, 0.05, "─" * 80, size=8, color=MUTED)

add_card(slide, 9, 1.5, 3.8, 4.8)

add_text(slide, 9.4, 1.8, 3.2, 0.4, "ИТОГО", size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 2.3, 3.2, 0.5, "+63 000 000", size=34, bold=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 2.9, 3.2, 0.3, "сум/мес доп. выручки", size=13, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 9.4, 3.6, 3.2, 0.3, "Стоимость Staffix:", size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 3.9, 3.2, 0.4, "260 000 сум/мес", size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text(slide, 9.4, 4.6, 3.2, 0.3, "Окупаемость:", size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 9.4, 4.9, 3.2, 0.5, "за 1 заказ", size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text(slide, 0.8, 6.6, 12, 0.4, "AI-сотрудник консультирует, продаёт и допродаёт. Вы собираете и отправляете.", size=16, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 0.8, 7.0, 12, 0.3, "staffix.io  •  14 дней бесплатно", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


# Save
output = r"c:\Users\anton\Documents\GITHUB REPOSITORIUM\staffix\Staffix_ROI_Slides.pptx"
prs.save(output)
print(f"Saved: {output}")
